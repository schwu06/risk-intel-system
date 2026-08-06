"""数据源文件解析、正文提取与可回溯文本切片。"""

from __future__ import annotations

import hashlib
import logging
import re
from dataclasses import dataclass, field, replace
from datetime import date, datetime
from functools import lru_cache
from html.parser import HTMLParser
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Optional
from urllib.parse import urlparse

import httpx
from docx import Document

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".txt", ".xlsx", ".docx", ".pptx", ".pdf"}
MAX_UPLOAD_BYTES = 25 * 1024 * 1024
MAX_EXTRACTED_CHARS = 120_000
CHUNK_TARGET_CHARS = 3_000
CHUNK_OVERLAP_CHARS = 120
PDF_OCR_DPI = 200
PDF_OCR_MAX_PAGES = 30
PDF_NATIVE_TEXT_THRESHOLD = 30


@dataclass(frozen=True)
class SourceChunk:
    """一个能够回到原始文档逻辑位置的文本切片。"""

    chunk_index: int
    text: str
    locator: str
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    cell_range: str | None = None
    row_range: str | None = None
    paragraph_index: int | None = None
    table_index: int | None = None
    table_row_index: int | None = None
    char_start: int | None = None
    char_end: int | None = None
    content_hash: str = ""


@dataclass(frozen=True)
class ParsedSource:
    extracted_text: str
    chunks: list[SourceChunk]
    parse_metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)
    is_truncated: bool = False
    used_ocr: bool = False


@dataclass(frozen=True)
class FetchedSource:
    """网址下载结果；raw_content专用于原始内容哈希。"""

    raw_content: bytes
    parsed: ParsedSource
    mime_type: str | None
    final_url: str


@dataclass(frozen=True)
class _Piece:
    text: str
    locator: str
    display_prefix: str = ""
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    cell_range: str | None = None
    row_range: str | None = None
    paragraph_index: int | None = None
    table_index: int | None = None
    table_row_index: int | None = None


def parse_uploaded_file(path: Path) -> str:
    """兼容旧调用：继续返回普通字符串。"""

    return parse_uploaded_file_with_chunks(path).extracted_text


def parse_uploaded_file_with_chunks(path: Path) -> ParsedSource:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        parsed = _parse_text(path.read_text(encoding="utf-8", errors="replace"))
    elif suffix == ".docx":
        parsed = _parse_docx_with_chunks(path)
    elif suffix == ".xlsx":
        parsed = _parse_xlsx_with_chunks(path)
    elif suffix == ".pptx":
        parsed = _parse_pptx_with_chunks(path)
    elif suffix == ".pdf":
        parsed = _parse_pdf_with_chunks(path)
    else:
        raise ValueError(f"不支持的文件类型: {suffix}")
    if not parsed.extracted_text.strip():
        raise ValueError("文件中未识别到可用于分析的正文")
    return parsed


def _sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _safe_cut_index(text: str, target: int, minimum: int = 1) -> int:
    """优先在自然边界切分，避免从数字、币种、百分号或年份中间断开。"""

    if len(text) <= target:
        return len(text)
    floor = max(minimum, target - 400)
    for index in range(target, floor - 1, -1):
        if _is_natural_boundary(text, index - 1):
            return index
    index = target
    protected = set("0123456789.,%％年月日万亿元美欧日人民币$¥￥+-")
    while index > minimum and (
        text[index - 1] in protected or (index < len(text) and text[index] in protected)
    ):
        index -= 1
    return index if index > minimum else target


def _split_piece(text: str) -> list[tuple[int, int, str]]:
    if len(text) <= CHUNK_TARGET_CHARS:
        return [(0, len(text), text)] if text else []
    result: list[tuple[int, int, str]] = []
    start = 0
    while start < len(text):
        remaining = text[start:]
        if len(remaining) <= CHUNK_TARGET_CHARS:
            end = len(text)
        else:
            end = start + _safe_cut_index(remaining, CHUNK_TARGET_CHARS)
        chunk = text[start:end].strip()
        if chunk:
            leading = len(text[start:end]) - len(text[start:end].lstrip())
            trailing = len(text[start:end].rstrip())
            actual_start = start + leading
            actual_end = start + trailing
            result.append((actual_start, actual_end, text[actual_start:actual_end]))
        if end >= len(text):
            break
        desired = max(start + 1, end - CHUNK_OVERLAP_CHARS)
        start = _safe_overlap_start(text, desired, end, minimum=start + 1)
    return result


def _safe_overlap_start(text: str, desired: int, end: int, *, minimum: int) -> int:
    for index in range(desired, min(end, desired + 200)):
        if _is_natural_boundary(text, index):
            return index + 1
    for index in range(desired - 1, max(minimum, desired - 200) - 1, -1):
        if _is_natural_boundary(text, index):
            return index + 1
    protected = set("0123456789.,%％年月日万亿元美欧日人民币$¥￥+-")
    index = desired
    while index > minimum and (
        text[index - 1] in protected or (index < len(text) and text[index] in protected)
    ):
        index -= 1
    return max(minimum, index)


def _is_natural_boundary(text: str, index: int) -> bool:
    char = text[index]
    if char not in "\n。！？；;.!?，,、 \t":
        return False
    if char in ".,，" and index > 0 and index + 1 < len(text):
        if text[index - 1].isdigit() and text[index + 1].isdigit():
            return False
    return True


def _assemble_pieces(
    pieces: list[_Piece],
    *,
    parse_metadata: dict[str, Any] | None = None,
    warnings: list[str] | None = None,
    used_ocr: bool = False,
) -> ParsedSource:
    rendered: list[str] = []
    chunks: list[SourceChunk] = []
    cursor = 0
    for piece in pieces:
        body = piece.text.strip()
        if not body:
            continue
        if rendered:
            rendered.append("\n\n")
            cursor += 2
        prefix = piece.display_prefix
        rendered.append(prefix)
        cursor += len(prefix)
        body_start = cursor
        rendered.append(body)
        cursor += len(body)
        subchunks = _split_piece(body)
        for part_index, (local_start, local_end, part_text) in enumerate(subchunks, start=1):
            locator = piece.locator
            if len(subchunks) > 1:
                locator = f"{locator}，片段{part_index}"
            chunks.append(
                SourceChunk(
                    chunk_index=len(chunks),
                    text=part_text,
                    locator=locator,
                    page_number=piece.page_number,
                    slide_number=piece.slide_number,
                    sheet_name=piece.sheet_name,
                    cell_range=piece.cell_range,
                    row_range=piece.row_range,
                    paragraph_index=piece.paragraph_index,
                    table_index=piece.table_index,
                    table_row_index=piece.table_row_index,
                    char_start=body_start + local_start,
                    char_end=body_start + local_end,
                    content_hash=_sha256_text(part_text),
                )
            )

    full_text = "".join(rendered).strip()
    warning_list = list(warnings or [])
    is_truncated = len(full_text) > MAX_EXTRACTED_CHARS
    if is_truncated:
        cut = _safe_cut_index(full_text, MAX_EXTRACTED_CHARS)
        retained = full_text[:cut].rstrip()
        full_text = retained + "\n...[内容已截断]"
        warning_list.append(
            f"解析文本超过{MAX_EXTRACTED_CHARS}字符，已在自然边界截断；尾部内容未进入系统。"
        )
        kept: list[SourceChunk] = []
        for chunk in chunks:
            if chunk.char_start is None or chunk.char_start >= len(retained):
                continue
            if chunk.char_end is not None and chunk.char_end > len(retained):
                shortened = retained[chunk.char_start :].strip()
                if shortened:
                    kept.append(
                        replace(
                            chunk,
                            text=shortened,
                            char_end=chunk.char_start + len(shortened),
                            content_hash=_sha256_text(shortened),
                        )
                    )
            else:
                kept.append(chunk)
        chunks = [replace(chunk, chunk_index=index) for index, chunk in enumerate(kept)]

    metadata = dict(parse_metadata or {})
    metadata["chunk_count"] = len(chunks)
    metadata["extracted_char_count"] = len(full_text)
    return ParsedSource(
        extracted_text=full_text,
        chunks=chunks,
        parse_metadata=metadata,
        warnings=warning_list,
        is_truncated=is_truncated,
        used_ocr=used_ocr,
    )


def _parse_text(text: str) -> ParsedSource:
    stripped = text.strip()
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", stripped) if part.strip()]
    if not paragraphs and stripped:
        paragraphs = [stripped]
    pieces = [
        _Piece(text=paragraph, locator=f"文本段落{index}", paragraph_index=index)
        for index, paragraph in enumerate(paragraphs, start=1)
    ]
    return _assemble_pieces(pieces, parse_metadata={"format": "txt"})


def parse_text_with_chunks(text: str, *, format_name: str = "txt") -> ParsedSource:
    """解析内存文本，供网址正文和搜索摘要登记复用。"""

    parsed = _parse_text(text)
    metadata = dict(parsed.parse_metadata)
    metadata["format"] = format_name
    return replace(parsed, parse_metadata=metadata)


def _parse_docx(path: Path) -> str:
    return _parse_docx_with_chunks(path).extracted_text


def _parse_docx_with_chunks(path: Path) -> ParsedSource:
    from docx.document import Document as _DocumentType
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    doc = Document(str(path))
    pieces: list[_Piece] = []
    paragraph_index = 0
    table_index = 0
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph_index += 1
            text = Paragraph(child, doc).text.strip()
            if text:
                pieces.append(
                    _Piece(
                        text=text,
                        locator=f"DOCX段落{paragraph_index}",
                        paragraph_index=paragraph_index,
                    )
                )
        elif isinstance(child, CT_Tbl):
            table_index += 1
            table = Table(child, doc if isinstance(doc, _DocumentType) else doc._parent)
            for row_index, row in enumerate(table.rows, start=1):
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    pieces.append(
                        _Piece(
                            text=" | ".join(cells),
                            locator=f"DOCX表格{table_index}第{row_index}行",
                            table_index=table_index,
                            table_row_index=row_index,
                        )
                    )
    warnings: list[str] = []
    if any(getattr(cell._tc, "tbl_lst", []) for table in doc.tables for row in table.rows for cell in row.cells):
        warnings.append("检测到嵌套表格；当前仅保证顶层段落和顶层表格的顺序及定位。")
    return _assemble_pieces(
        pieces,
        parse_metadata={
            "format": "docx",
            "paragraph_count": paragraph_index,
            "table_count": table_index,
        },
        warnings=warnings,
    )


def _xlsx_scalar(value: Any) -> str:
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, date):
        return value.isoformat()
    return str(value)


def _xlsx_cell_text(value_cell, formula_cell) -> str:
    value = value_cell.value
    formula = formula_cell.value
    if value is None and formula is None:
        return ""
    rendered = _xlsx_scalar(value if value is not None else formula)
    if isinstance(formula, str) and formula.startswith("="):
        rendered = f"{rendered} [公式: {formula}]"
    number_format = str(formula_cell.number_format or "General")
    if number_format != "General" and isinstance(value, (int, float)):
        rendered = f"{rendered} [格式: {number_format}]"
    return rendered


def _parse_xlsx(path: Path) -> str:
    return _parse_xlsx_with_chunks(path).extracted_text


def _parse_xlsx_with_chunks(path: Path) -> ParsedSource:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise RuntimeError("缺少 openpyxl 依赖，无法解析 xlsx") from exc

    values_wb = load_workbook(str(path), read_only=True, data_only=True)
    formulas_wb = load_workbook(str(path), read_only=True, data_only=False)
    pieces: list[_Piece] = []
    sheet_count = len(values_wb.sheetnames)
    try:
        for values_ws in values_wb.worksheets:
            formulas_ws = formulas_wb[values_ws.title]
            first_row_in_sheet = True
            row_pairs = zip(values_ws.iter_rows(), formulas_ws.iter_rows())
            for row_index, (value_row, formula_row) in enumerate(row_pairs, start=1):
                rendered_cells: list[tuple[int, str]] = []
                for column_index, value_cell in enumerate(value_row, start=1):
                    text = _xlsx_cell_text(value_cell, formula_row[column_index - 1])
                    if text:
                        rendered_cells.append((column_index, text))
                if not rendered_cells:
                    continue
                first_col = rendered_cells[0][0]
                last_col = rendered_cells[-1][0]
                cell_range = (
                    f"{get_column_letter(first_col)}{row_index}:"
                    f"{get_column_letter(last_col)}{row_index}"
                )
                prefix = f"[工作表: {values_ws.title}]\n" if first_row_in_sheet else ""
                first_row_in_sheet = False
                pieces.append(
                    _Piece(
                        text=" | ".join(text for _, text in rendered_cells),
                        display_prefix=prefix,
                        locator=f"工作表“{values_ws.title}”第{row_index}行（{cell_range}）",
                        sheet_name=values_ws.title,
                        cell_range=cell_range,
                        row_range=str(row_index),
                    )
                )
    finally:
        values_wb.close()
        formulas_wb.close()
    return _assemble_pieces(
        pieces,
        parse_metadata={"format": "xlsx", "sheet_count": sheet_count},
    )


def _parse_pptx(path: Path) -> str:
    return _parse_pptx_with_chunks(path).extracted_text


def _parse_pptx_with_chunks(path: Path) -> ParsedSource:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx 依赖，无法解析 pptx") from exc

    presentation = Presentation(str(path))
    pieces: list[_Piece] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            slide_parts.extend(_pptx_shape_text(shape))
        if slide_parts:
            pieces.append(
                _Piece(
                    text="\n".join(slide_parts),
                    display_prefix=f"[幻灯片 {slide_index}]\n",
                    locator=f"PPTX幻灯片{slide_index}正文",
                    slide_number=slide_index,
                )
            )
        try:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            pieces.append(
                _Piece(
                    text=notes,
                    display_prefix="[备注]\n",
                    locator=f"PPTX幻灯片{slide_index}备注",
                    slide_number=slide_index,
                )
            )
    return _assemble_pieces(
        pieces,
        parse_metadata={"format": "pptx", "slide_count": len(presentation.slides)},
    )


def _pptx_shape_text(shape) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    for child in getattr(shape, "shapes", ()):
        parts.extend(_pptx_shape_text(child))
    return parts


def _parse_pdf(path: Path) -> str:
    return _parse_pdf_with_chunks(path).extracted_text


def _parse_pdf_with_chunks(path: Path) -> ParsedSource:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("缺少 pypdf 依赖，无法解析 pdf") from exc

    reader = PdfReader(str(path))
    native_pages = [(page.extract_text() or "").strip() for page in reader.pages]
    ocr_indexes = [
        index
        for index, text in enumerate(native_pages)
        if len(re.sub(r"\s+", "", text)) < PDF_NATIVE_TEXT_THRESHOLD
    ]
    if len(ocr_indexes) > PDF_OCR_MAX_PAGES:
        raise ValueError(f"扫描 PDF 需要 OCR 的页数超过上限 {PDF_OCR_MAX_PAGES} 页")

    ocr_pages = _ocr_pdf_pages(path, ocr_indexes) if ocr_indexes else {}
    pieces: list[_Piece] = []
    for index, native_text in enumerate(native_pages):
        page_number = index + 1
        text = ocr_pages.get(index) or native_text
        if text.strip():
            used_page_ocr = index in ocr_pages
            marker = " OCR" if used_page_ocr else ""
            pieces.append(
                _Piece(
                    text=text,
                    display_prefix=f"[PDF 第 {page_number} 页{marker}]\n",
                    locator=f"PDF第{page_number}页",
                    page_number=page_number,
                )
            )
    warnings = []
    if ocr_pages:
        pages = "、".join(str(index + 1) for index in sorted(ocr_pages))
        warnings.append(f"PDF第{pages}页使用OCR文本，数字和符号应结合原页复核。")
    missed_ocr = [index + 1 for index in ocr_indexes if index not in ocr_pages]
    if missed_ocr:
        warnings.append(f"PDF第{'、'.join(map(str, missed_ocr))}页未识别到足够正文。")
    return _assemble_pieces(
        pieces,
        parse_metadata={
            "format": "pdf",
            "page_count": len(native_pages),
            "ocr_pages": [index + 1 for index in sorted(ocr_pages)],
        },
        warnings=warnings,
        used_ocr=bool(ocr_pages),
    )


@lru_cache(maxsize=1)
def _get_ocr_engine():
    try:
        from rapidocr import RapidOCR
    except ImportError as exc:
        raise RuntimeError("缺少 rapidocr 依赖，无法识别扫描 PDF") from exc
    return RapidOCR()


def _ocr_pdf_pages(path: Path, page_indexes: list[int]) -> dict[int, str]:
    try:
        import pymupdf
    except ImportError as exc:
        raise RuntimeError("缺少 PyMuPDF 依赖，无法渲染扫描 PDF") from exc

    if not page_indexes:
        return {}
    engine = _get_ocr_engine()
    recognized: dict[int, str] = {}
    document = pymupdf.open(str(path))
    try:
        for index in page_indexes:
            page = document[index]
            pixmap = page.get_pixmap(
                dpi=PDF_OCR_DPI,
                colorspace=pymupdf.csRGB,
                alpha=False,
            )
            result = engine(pixmap.tobytes("png"))
            texts = getattr(result, "txts", None) or ()
            text = "\n".join(str(item).strip() for item in texts if str(item).strip())
            if text:
                recognized[index] = text
    except Exception as exc:
        raise RuntimeError(f"扫描 PDF OCR 失败: {exc}") from exc
    finally:
        document.close()
    return recognized


class _ReadableHTMLParser(HTMLParser):
    _BLOCK_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "blockquote", "td", "th"}
    _CONTAINER_TAGS = {"article", "main", "section", "div"}
    _IGNORED_TAGS = {"script", "style", "nav", "header", "footer", "aside", "noscript", "svg"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: list[str] = []
        self._ignored_depth = 0
        self._active_tag: str | None = None
        self._buffer: list[str] = []
        self._container_stack: list[list[str]] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        tag = tag.lower()
        if tag in self._IGNORED_TAGS:
            self._ignored_depth += 1
            return
        if self._ignored_depth:
            return
        if tag in self._BLOCK_TAGS:
            self._flush_container_buffer()
            self._flush()
            self._active_tag = tag
        elif tag in self._CONTAINER_TAGS:
            self._flush_container_buffer()
            self._container_stack.append([])

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in self._IGNORED_TAGS:
            self._ignored_depth = max(0, self._ignored_depth - 1)
            return
        if not self._ignored_depth:
            if tag == self._active_tag:
                self._flush()
            elif tag in self._CONTAINER_TAGS and self._container_stack:
                self._flush_container_buffer()
                self._container_stack.pop()

    def handle_data(self, data: str) -> None:
        if self._ignored_depth:
            return
        if self._active_tag:
            self._buffer.append(data)
        elif self._container_stack:
            self._container_stack[-1].append(data)

    def close(self) -> None:
        super().close()
        self._flush()
        while self._container_stack:
            self._flush_container_buffer()
            self._container_stack.pop()

    def _flush(self) -> None:
        text = re.sub(r"\s+", " ", " ".join(self._buffer)).strip()
        if text:
            self.blocks.append(text)
        self._buffer = []
        self._active_tag = None

    def _flush_container_buffer(self) -> None:
        if not self._container_stack:
            return
        buffer = self._container_stack[-1]
        text = re.sub(r"\s+", " ", " ".join(buffer)).strip()
        if text:
            self.blocks.append(text)
        buffer.clear()


def parse_html_with_chunks(html_text: str) -> ParsedSource:
    parser = _ReadableHTMLParser()
    parser.feed(html_text)
    parser.close()
    pieces = [
        _Piece(text=text, locator=f"网页段落{index}", paragraph_index=index)
        for index, text in enumerate(parser.blocks, start=1)
    ]
    return _assemble_pieces(pieces, parse_metadata={"format": "html", "paragraph_count": len(pieces)})


def fetch_url_source(url: str, timeout: int = 20) -> FetchedSource:
    parsed_url = urlparse(url)
    origin = f"{parsed_url.scheme}://{parsed_url.netloc}/"
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/138.0.0.0 Safari/537.36"
        ),
        "Accept": "application/pdf,text/html,application/xhtml+xml,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        "Referer": origin,
        "Cache-Control": "no-cache",
    }
    with httpx.Client(timeout=timeout, follow_redirects=True) as client:
        with client.stream("GET", url, headers=headers) as resp:
            if resp.status_code == 403:
                raise ValueError(
                    "PDF 所在网站拒绝服务器下载（HTTP 403）。"
                    "该网站可能要求登录、Cookie 或防盗链验证，请先下载 PDF，再使用文件上传。"
                )
            resp.raise_for_status()
            content_type = (resp.headers.get("content-type") or "").split(";", 1)[0].strip().lower()
            content = bytearray()
            for block in resp.iter_bytes():
                content.extend(block)
                if len(content) > MAX_UPLOAD_BYTES:
                    raise ValueError("网址内容超过 25 MB 上限")
            encoding = resp.charset_encoding or "utf-8"
            final_url = str(resp.url)

    raw = bytes(content)
    path_is_pdf = urlparse(final_url).path.lower().endswith(".pdf")
    if content_type == "application/pdf" or path_is_pdf or raw.startswith(b"%PDF-"):
        with TemporaryDirectory(prefix="risk-intel-url-pdf-") as tmp:
            pdf_path = Path(tmp) / "source.pdf"
            pdf_path.write_bytes(raw)
            parsed = parse_uploaded_file_with_chunks(pdf_path)
        mime_type = "application/pdf"
    else:
        text = raw.decode(encoding, errors="replace")
        if "html" in content_type or "<html" in text[:500].lower():
            parsed = parse_html_with_chunks(text)
            mime_type = content_type or "text/html"
        else:
            parsed = _parse_text(text)
            mime_type = content_type or "text/plain"
        if not parsed.extracted_text.strip():
            raise ValueError("网址中未识别到可用于分析的正文")
    return FetchedSource(
        raw_content=raw,
        parsed=parsed,
        mime_type=mime_type,
        final_url=final_url,
    )


def fetch_url_text(url: str, timeout: int = 20) -> str:
    """兼容旧调用：继续返回网址正文字符串。"""

    return fetch_url_source(url, timeout=timeout).parsed.extracted_text


def _strip_html(html_text: str) -> str:
    """兼容旧内部调用；现在保留正文段落而不是压成单行。"""

    return parse_html_with_chunks(html_text).extracted_text


def truncate_text(text: str, max_chars: int = MAX_EXTRACTED_CHARS) -> str:
    if len(text) <= max_chars:
        return text
    cut = _safe_cut_index(text, max_chars)
    return text[:cut].rstrip() + "\n...[内容已截断]"


def detect_extension(filename: str) -> Optional[str]:
    ext = Path(filename).suffix.lower()
    return ext if ext in SUPPORTED_EXTENSIONS else None
