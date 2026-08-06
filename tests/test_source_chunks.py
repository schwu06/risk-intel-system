from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.services.data_source_parser import (
    MAX_EXTRACTED_CHARS,
    parse_html_with_chunks,
    parse_text_with_chunks,
    parse_uploaded_file_with_chunks,
)


class _PdfPage:
    def __init__(self, text: str) -> None:
        self.text = text

    def extract_text(self) -> str:
        return self.text


class _PdfReader:
    def __init__(self, _path: str) -> None:
        self.pages = [_PdfPage("第一页短文本"), _PdfPage("第二页短文本")]


class SourceChunkParserTests(unittest.TestCase):
    def test_pdf_chunks_keep_page_and_ocr_location(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "source.pdf"
            path.write_bytes(b"%PDF-test")
            with (
                patch("pypdf.PdfReader", _PdfReader),
                patch(
                    "app.services.data_source_parser._ocr_pdf_pages",
                    return_value={0: "OCR第一页金额100亿元", 1: "OCR第二页增长12.5%"},
                ),
            ):
                parsed = parse_uploaded_file_with_chunks(path)
        self.assertTrue(parsed.used_ocr)
        self.assertEqual([chunk.page_number for chunk in parsed.chunks], [1, 2])
        self.assertEqual(parsed.chunks[0].locator, "PDF第1页")
        self.assertIn("第1、2页使用OCR", parsed.warnings[0])

    def test_docx_preserves_top_level_paragraph_table_order(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "source.docx"
            doc = Document()
            doc.add_paragraph("表格之前")
            table = doc.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "收入"
            table.cell(0, 1).text = "100亿元"
            doc.add_paragraph("表格之后")
            doc.save(path)
            parsed = parse_uploaded_file_with_chunks(path)
        self.assertLess(parsed.extracted_text.index("表格之前"), parsed.extracted_text.index("收入"))
        self.assertLess(parsed.extracted_text.index("收入"), parsed.extracted_text.index("表格之后"))
        self.assertEqual(parsed.chunks[0].paragraph_index, 1)
        self.assertEqual(parsed.chunks[1].table_index, 1)
        self.assertEqual(parsed.chunks[1].table_row_index, 1)
        self.assertEqual(parsed.chunks[2].paragraph_index, 2)

    def test_xlsx_keeps_sheet_cell_range_value_format_and_formula(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "source.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = "财务"
            ws["A1"] = "增长率"
            ws["B1"] = 0.125
            ws["B1"].number_format = "0.0%"
            ws["A2"] = "合计"
            ws["B2"] = "=B1*2"
            wb.save(path)
            wb.close()
            parsed = parse_uploaded_file_with_chunks(path)
        self.assertEqual(parsed.parse_metadata["sheet_count"], 1)
        self.assertEqual(parsed.chunks[0].sheet_name, "财务")
        self.assertEqual(parsed.chunks[0].cell_range, "A1:B1")
        self.assertIn("0.125", parsed.chunks[0].text)
        self.assertIn("格式: 0.0%", parsed.chunks[0].text)
        self.assertIn("公式: =B1*2", parsed.chunks[1].text)

    def test_pptx_distinguishes_slide_body(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "source.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "行业风险"
            slide.placeholders[1].text = "正文内容"
            presentation.save(path)
            parsed = parse_uploaded_file_with_chunks(path)
        self.assertEqual(parsed.parse_metadata["slide_count"], 1)
        self.assertEqual(parsed.chunks[0].slide_number, 1)
        self.assertIn("正文", parsed.chunks[0].locator)

    def test_html_keeps_paragraphs_and_removes_navigation(self) -> None:
        parsed = parse_html_with_chunks(
            "<html><nav>菜单</nav><h1>标题</h1><p>第一段100亿元</p>"
            "<script>bad()</script><p>第二段12.5%</p></html>"
        )
        self.assertNotIn("菜单", parsed.extracted_text)
        self.assertNotIn("bad", parsed.extracted_text)
        self.assertEqual([chunk.paragraph_index for chunk in parsed.chunks], [1, 2, 3])

    def test_html_extracts_div_and_article_body_without_duplication(self) -> None:
        parsed = parse_html_with_chunks(
            "<html><body><article><h1>行业标题</h1>"
            "<div>第一段收入100亿元</div>"
            "<section><div>第二段增长12.5%</div></section>"
            "</article></body></html>"
        )
        self.assertEqual(
            [chunk.text for chunk in parsed.chunks],
            ["行业标题", "第一段收入100亿元", "第二段增长12.5%"],
        )
        self.assertEqual([chunk.paragraph_index for chunk in parsed.chunks], [1, 2, 3])

    def test_long_text_is_explicitly_marked_truncated(self) -> None:
        text = ("金额100亿元，增长12.5%。\n\n" * 7000) + "尾部"
        self.assertGreater(len(text), MAX_EXTRACTED_CHARS)
        parsed = parse_text_with_chunks(text)
        self.assertTrue(parsed.is_truncated)
        self.assertTrue(parsed.extracted_text.endswith("...[内容已截断]"))
        self.assertTrue(parsed.warnings)
        self.assertTrue(all((c.char_end or 0) <= len(parsed.extracted_text) for c in parsed.chunks))

    def test_chunk_boundary_does_not_split_number_currency_and_unit(self) -> None:
        token = "人民币123456.78亿元"
        parsed = parse_text_with_chunks(("甲" * 2950) + token + ("乙" * 3500))
        touching = [chunk.text for chunk in parsed.chunks if "123456" in chunk.text]
        self.assertTrue(touching)
        self.assertTrue(all(token in text for text in touching))


if __name__ == "__main__":
    unittest.main()
